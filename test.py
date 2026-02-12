from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os
from datetime import date

# ── Configuration ──────────────────────────────────────────────────────
TEMPLATE_PATH = None  # Set to "template.pptx" if you have the file
OUTPUT_FILE = "System_Team_Leadership_Update_2026-02-11.pptx"

# ESDC Brand Colors (from template)
ESDC_TEAL       = RGBColor(0x26, 0x97, 0x8F)   # Primary teal
ESDC_TEAL_DARK  = RGBColor(0x1A, 0x6E, 0x68)   # Dark teal
ESDC_TEAL_LIGHT = RGBColor(0x8C, 0xC9, 0xC3)   # Light teal
ESDC_CORAL      = RGBColor(0xE8, 0x52, 0x52)    # Accent coral/red
ESDC_CORAL_LIGHT= RGBColor(0xF2, 0x9E, 0x9E)   # Light coral
ESDC_DARK       = RGBColor(0x33, 0x33, 0x33)    # Body text
ESDC_GREY       = RGBColor(0x66, 0x66, 0x66)    # Secondary text
ESDC_LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)   # Background fills
WHITE           = RGBColor(0xFF, 0xFF, 0xFF)
BLACK           = RGBColor(0x00, 0x00, 0x00)

# Status colors
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
AMBER  = RGBColor(0xF5, 0x7F, 0x17)
RED_ST = RGBColor(0xC6, 0x28, 0x28)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Helper Functions ───────────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a shape with optional fill and border."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=12,
                font_color=ESDC_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", word_wrap=True):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_frame(slide, left, top, width, height, items,
                     font_size=14, font_color=ESDC_DARK, bullet_color=ESDC_TEAL,
                     spacing_after=Pt(6), font_name="Calibri", bold_items=None):
    """Add a text frame with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_after = spacing_after
        p.level = 0
        if bold_items and i in bold_items:
            p.font.bold = True
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '▸'})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = buClr.makeelement(qn('a:srgbClr'), {'val': f'{bullet_color[0]:02X}{bullet_color[1]:02X}{bullet_color[2]:02X}'})  # noqa
        buClr.append(srgb)
        pPr.append(buClr)
        pPr.append(buChar)
    return txBox

def add_esdc_header_bar(slide):
    """Add the teal tri-color header bar (ESDC branding)."""
    bar_top = Inches(0.55)
    bar_h = Inches(0.06)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), bar_top, Inches(8.0), bar_h, ESDC_TEAL)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(8.6), bar_top, Inches(2.0), bar_h, ESDC_TEAL_LIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(10.7), bar_top, Inches(2.1), bar_h, ESDC_LIGHT_GREY)

def add_slide_number(slide, number):
    """Add slide number bottom-right."""
    add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35),
                str(number), font_size=10, font_color=ESDC_GREY,
                alignment=PP_ALIGN.RIGHT)

def add_footer_bar(slide):
    """Add the colorful mosaic footer bar."""
    bar_top = Inches(7.1)
    bar_h = Inches(0.18)
    colors = [ESDC_TEAL, ESDC_TEAL_LIGHT, ESDC_CORAL, ESDC_CORAL_LIGHT,
              ESDC_TEAL_DARK, RGBColor(0xB0, 0xDB, 0xD6)]
    seg_w = Inches(13.333 / 30)
    for i in range(30):
        c = colors[i % len(colors)]
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  Emu(int(i * seg_w)), bar_top, seg_w, bar_h, c)

def add_callout_box(slide, left, top, width, height, text,
                    fill=ESDC_TEAL, text_color=WHITE, font_size=13, bold=True):
    """Add a colored callout / highlight box."""
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    box.text_frame.paragraphs[0].space_before = Pt(4)
    return box

def add_status_badge(slide, left, top, text, color):
    """Small status indicator badge."""
    w, h = Inches(1.5), Inches(0.35)
    box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h, color)
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return box

def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text

def add_section_divider_accent(slide):
    """Add the geometric accent in upper-right for section headers."""
    # Simplified geometric accent using overlapping shapes
    base_left = Inches(11.0)
    base_top = Inches(0.0)
    colors = [ESDC_TEAL_LIGHT, ESDC_CORAL_LIGHT, RGBColor(0xD4, 0xEC, 0xE9)]
    for i, c in enumerate(colors):
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  base_left + Inches(i * 0.4), base_top + Inches(i * 0.3),
                  Inches(1.2), Inches(1.2), c)


# ── Slide Builders ─────────────────────────────────────────────────────

def build_title_slide(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, WHITE)

    # Teal accent bar across top
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ESDC_TEAL)

    # ESDC header text
    add_textbox(slide, Inches(1.0), Inches(0.3), Inches(6), Inches(0.5),
                "Employment and Social Development Canada", font_size=11,
                font_color=ESDC_DARK, bold=False)

    # Tri-color bar
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.8), Inches(8.5), Inches(0.06), ESDC_TEAL)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(8.6), Inches(0.8), Inches(2.3), Inches(0.06), ESDC_TEAL_LIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(11.0), Inches(0.8), Inches(2.3), Inches(0.06), ESDC_LIGHT_GREY)

    # Title
    add_textbox(slide, Inches(3.5), Inches(2.8), Inches(7.5), Inches(1.2),
                "System Team Operating Model,\nEnablement Activation &\nQuality Dashboard",
                font_size=30, font_color=ESDC_DARK, bold=True, alignment=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, Inches(3.5), Inches(4.3), Inches(7.5), Inches(0.5),
                "Leadership Update  |  EI on BDM Program",
                font_size=16, font_color=ESDC_GREY, alignment=PP_ALIGN.LEFT)

    # Date
    add_textbox(slide, Inches(3.5), Inches(4.9), Inches(4), Inches(0.4),
                f"February 11, 2026", font_size=14, font_color=ESDC_TEAL)

    # Left geometric accent (simplified)
    colors_geo = [ESDC_TEAL, ESDC_TEAL_LIGHT, ESDC_CORAL, ESDC_CORAL_LIGHT]
    for i in range(6):
        c = colors_geo[i % len(colors_geo)]
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  Inches(0.3 + i * 0.35), Inches(1.5 + i * 0.55),
                  Inches(1.0), Inches(1.0), c)

    add_speaker_notes(slide,
        "OPENING:\n"
        "Good morning/afternoon. Today I'm providing a leadership update on three interconnected initiatives "
        "that are critical to our delivery maturity.\n\n"
        "We are closing two core enablers: the System Team Charter as our ratified operating model, and "
        "Quality Dashboards as our trusted signal for release health. In parallel, we're activating "
        "Accessibility and Automation enablement through ARTs and STEs.\n\n"
        "TRANSITION: Let me start with what you need to know at a glance.")


def build_exec_summary(prs):
    """Slide 2: Executive Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    n = 2

    # Title
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "EI on BDM  —  Leadership Update", font_size=14,
                font_color=ESDC_GREY, bold=False)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.5),
                "Executive Summary: Three Enablers Converging to Unlock Consistent Delivery",
                font_size=22, font_color=ESDC_DARK, bold=True)

    # Three column callout boxes
    col_w = Inches(3.8)
    col_gap = Inches(0.25)
    top = Inches(1.6)
    box_h = Inches(0.55)

    headers = [
        "1  System Team Charter",
        "2  Quality Dashboard & KPIs",
        "3  Enablement Activation"
    ]
    statuses = ["On Track", "Green (1 dependency)", "Activation Week of Feb 12"]
    status_colors = [GREEN, GREEN, ESDC_TEAL]

    for i, (hdr, st, sc) in enumerate(zip(headers, statuses, status_colors)):
        x = Inches(0.5) + (col_w + col_gap) * i
        add_callout_box(slide, x, top, col_w, box_h, hdr,
                        fill=ESDC_TEAL, text_color=WHITE, font_size=14)
        add_status_badge(slide, x + Inches(1.15), top + Inches(0.65), st, sc)

    # Key content under each column
    col_items = [
        [
            "Ratify operating model: purpose, scope, roles, decision rights",
            "3-step sequence: Align → Activate → Embed",
            "4 auditable milestones with clear owners",
            "Reduces ambiguity, speeds cross-team decisions"
        ],
        [
            "Baseline KPIs locked; RCA KPI in progress",
            "Power BI dependency managed (mitigation in place)",
            "Becomes single source of truth for CI signals",
            "Shifts decisions from opinion-based to signal-based"
        ],
        [
            "Accessibility & Automation via ARTs/STEs",
            "Internal greenlight: Feb 12",
            "ART/STE kickoff: Feb 18",
            "Product team sessions: week of Feb 23"
        ]
    ]

    for i, items in enumerate(col_items):
        x = Inches(0.5) + (col_w + col_gap) * i
        add_bullet_frame(slide, x + Inches(0.1), Inches(2.65), col_w - Inches(0.2), Inches(3.5),
                         items, font_size=12, bullet_color=ESDC_TEAL, spacing_after=Pt(8))

    # Bottom callout — leadership asks
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.95),
              Inches(12.3), Inches(0.05), ESDC_CORAL)

    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(2), Inches(0.35),
                "3 DECISIONS NEEDED", font_size=13, font_color=ESDC_CORAL, bold=True)

    asks = [
        "❶ Confirm sponsorship for Charter activation & CI alignment with Chris and Don",
        "❷ Confirm Power BI capacity or approve short-term QE backfill for KPI automation",
        "❸ Confirm Quality Dashboard as default view for weekly release health & executive updates"
    ]
    add_bullet_frame(slide, Inches(0.5), Inches(6.4), Inches(12), Inches(1.0),
                     asks, font_size=12, font_color=ESDC_DARK, bullet_color=ESDC_CORAL,
                     spacing_after=Pt(3))

    add_slide_number(slide, n)

    add_speaker_notes(slide,
        "KEY MESSAGE:\n"
        "We're converging three related workstreams that together create a flywheel: "
        "a ratified operating model tells us HOW we work, the quality dashboard tells us "
        "HOW WELL we work, and enablement activation puts the tools in teams' hands.\n\n"
        "EMPHASIS:\n"
        "- All three are green or on track. The one managed risk is Power BI activation capacity.\n"
        "- I have three specific asks at the bottom — I'll come back to these at the end.\n\n"
        "ANTICIPATED QUESTION: 'Why do these need to happen together?'\n"
        "ANSWER: The Charter defines the interface; the dashboard provides accountability; "
        "enablement activates the teams. Without all three, we get process without measurement, "
        "or measurement without adoption.\n\n"
        "TRANSITION: Let me walk you through the agenda.")


def build_agenda(prs):
    """Slide 3: Agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    add_section_divider_accent(slide)
    n = 3

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "EI on BDM  —  Leadership Update", font_size=14,
                font_color=ESDC_GREY)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(10), Inches(0.5),
                "Agenda", font_size=26, font_color=ESDC_DARK, bold=True)

    items = [
        ("01", "What We Are Driving", "Strategic context and convergence of three enablers"),
        ("02", "System Team Charter", "Operating model finalization — Align, Activate, Embed"),
        ("03", "Quality Dashboard & KPIs", "Trusted signal for release health and shift-left progress"),
        ("04", "Enablement Activation", "Accessibility & Automation operationalized through ARTs/STEs"),
        ("05", "Timeline & Milestones", "Key dates and auditable closure plan"),
        ("06", "Leadership Asks", "Three decisions needed to maintain momentum"),
    ]

    for i, (num, title, desc) in enumerate(items):
        y = Inches(1.65) + Inches(i * 0.85)
        # Number circle
        circ = add_shape(slide, MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.55), Inches(0.55), ESDC_TEAL)
        p = circ.text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        circ.text_frame.paragraphs[0].space_before = Pt(2)

        add_textbox(slide, Inches(1.5), y, Inches(5), Inches(0.35),
                    title, font_size=17, font_color=ESDC_DARK, bold=True)
        add_textbox(slide, Inches(1.5), y + Inches(0.32), Inches(8), Inches(0.3),
                    desc, font_size=13, font_color=ESDC_GREY)

    # Estimated time
    add_callout_box(slide, Inches(9.5), Inches(1.65), Inches(2.5), Inches(0.45),
                    "~15 min + Q&A", fill=ESDC_TEAL_LIGHT, text_color=ESDC_TEAL_DARK, font_size=13)

    add_slide_number(slide, n)
    add_speaker_notes(slide,
        "TALKING POINT:\n"
        "Here's our roadmap for today. I'll spend the bulk of the time on the three workstreams "
        "and close with the specific decisions I need from you.\n\n"
        "TRANSITION: Let's start with the strategic context — what we're driving and why.")


def build_strategic_context(prs):
    """Slide 4: What We Are Driving."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    n = 4

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "EI on BDM  —  Leadership Update", font_size=14, font_color=ESDC_GREY)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.5),
                "Two Core Enablers Unlock Consistent Delivery and Measurable Improvement",
                font_size=22, font_color=ESDC_DARK, bold=True)

    # Left box — Enabler 1
    box1 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                     Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5), ESDC_LIGHT_GREY)
    box1.line.color.rgb = ESDC_TEAL
    box1.line.width = Pt(1.5)

    add_callout_box(slide, Inches(0.7), Inches(1.75), Inches(3.2), Inches(0.45),
                    "ENABLER 1: Operating Model", fill=ESDC_TEAL)

    add_bullet_frame(slide, Inches(0.8), Inches(2.35), Inches(5.2), Inches(1.6), [
        "System Team Charter as the ratified execution baseline",
        "Defines purpose, scope, roles, decision rights, cadences",
        "Makes ownership and escalation paths explicit",
    ], font_size=13, bullet_color=ESDC_TEAL, spacing_after=Pt(6))

    # Right box — Enabler 2
    box2 = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                     Inches(6.7), Inches(1.6), Inches(5.8), Inches(2.5), ESDC_LIGHT_GREY)
    box2.line.color.rgb = ESDC_CORAL
    box2.line.width = Pt(1.5)

    add_callout_box(slide, Inches(6.9), Inches(1.75), Inches(3.8), Inches(0.45),
                    "ENABLER 2: Quality Signal", fill=ESDC_CORAL)

    add_bullet_frame(slide, Inches(7.0), Inches(2.35), Inches(5.2), Inches(1.6), [
        "Quality Dashboards & KPIs as trusted release health signal",
        "Tracks shift-left progress and intervention points",
        "Enables signal-based (not opinion-based) decisions",
    ], font_size=13, bullet_color=ESDC_CORAL, spacing_after=Pt(6))

    # Arrow / connector area
    add_shape(slide, MSO_SHAPE.DOWN_ARROW, Inches(6.0), Inches(4.4),
              Inches(1.3), Inches(0.6), ESDC_TEAL)

    # Bottom box — Convergence
    conv_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(2.0), Inches(5.2), Inches(9.3), Inches(1.6), ESDC_TEAL_DARK)

    txf = conv_box.text_frame
    txf.word_wrap = True
    p = txf.paragraphs[0]
    p.text = "CONVERGENCE: Enablement Activation"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(4)
    p2 = txf.add_paragraph()
    p2.text = ("Accessibility and Automation enablement is operationalized through ARTs and STEs, "
               "then embedded at the Product Team level using the dashboard and metrics. "
               "This directly supports the Charter intent of a consistent interface, clear ownership, "
               "and closed-loop feedback.")
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0xD4, 0xEC, 0xE9)
    p2.alignment = PP_ALIGN.LEFT

    add_slide_number(slide, n)
    add_speaker_notes(slide,
        "KEY MESSAGE:\n"
        "These aren't three separate projects — they're an interlocking system.\n\n"
        "The Charter tells us HOW we operate. The Dashboard tells us HOW WELL we're doing. "
        "Enablement Activation puts the tools and practices into product teams' hands.\n\n"
        "Without the Charter, enablement has no consistent interface. Without the Dashboard, "
        "we can't measure adoption. Without activation, the Charter is just a document on SharePoint.\n\n"
        "TRANSITION: Let's go deeper into each one, starting with the System Team Charter.")


def build_charter_section_header(prs):
    """Slide 5: Section header — System Team Charter."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_divider_accent(slide)
    n = 5

    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(8), Inches(0.6),
                "System Team Charter", font_size=32, font_color=ESDC_DARK, bold=True)

    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(8), Inches(0.5),
                "Finalize the operating model and activate it as the execution baseline",
                font_size=16, font_color=ESDC_GREY)

    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(3.0), Inches(0.05), ESDC_TEAL)

    add_slide_number(slide, n)
    add_speaker_notes(slide, "TRANSITION: Section 1 — System Team Charter.")


def build_charter_align_activate_embed(prs):
    """Slide 6: Charter — 3-step enablement sequence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    n = 6

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "System Team Charter", font_size=14, font_color=ESDC_GREY)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.5),
                "Three-Step Sequence Moves the Charter from Document to Operating Rhythm",
                font_size=22, font_color=ESDC_DARK, bold=True)

    phases = [
        {
            "title": "ALIGN",
            "color": ESDC_TEAL_LIGHT,
            "items": [
                "Finalize Charter v1.0: purpose, scope,\nR&Rs, decision rights, cadences",
                "Confirm changes vs. current ways of working",
            ],
            "output": "Charter ready for ratification + one-page change summary"
        },
        {
            "title": "ACTIVATE",
            "color": ESDC_TEAL,
            "items": [
                "Run System Team enablement session:\n  • QPA insights and execution impact\n  • Charter walkthrough and escalation paths",
                "Confirm adoption mechanics: cadence,\nartifacts, owners",
            ],
            "output": "Ratified Charter, decision log updated, action items assigned"
        },
        {
            "title": "EMBED",
            "color": ESDC_TEAL_DARK,
            "items": [
                "Operationalize for next 2 sprints",
                "Use Charter in ceremonies: standups,\nrefinement, release pulse",
                "Enforce lightweight governance: decision\nauthority & DoD for System Team items",
            ],
            "output": "Adoption evidence: cleaner dependency hygiene, fewer late integration surprises"
        }
    ]

    for i, phase in enumerate(phases):
        x = Inches(0.5) + Inches(i * 4.2)
        # Phase header
        add_callout_box(slide, x, Inches(1.55), Inches(3.8), Inches(0.5),
                        phase["title"], fill=phase["color"], font_size=16)

        # Content
        y_start = Inches(2.3)
        add_bullet_frame(slide, x + Inches(0.1), y_start, Inches(3.6), Inches(2.8),
                         phase["items"], font_size=12, bullet_color=phase["color"],
                         spacing_after=Pt(8))

        # Output box
        out_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                            x, Inches(5.3), Inches(3.8), Inches(0.9), ESDC_LIGHT_GREY)
        out_box.line.color.rgb = phase["color"]
        out_box.line.width = Pt(1)

        add_textbox(slide, x + Inches(0.1), Inches(5.35), Inches(3.6), Inches(0.2),
                    "OUTPUT:", font_size=10, font_color=phase["color"], bold=True)
        add_textbox(slide, x + Inches(0.1), Inches(5.55), Inches(3.6), Inches(0.6),
                    phase["output"], font_size=11, font_color=ESDC_DARK)

        # Arrow connectors between phases
        if i < 2:
            add_shape(slide, MSO_SHAPE.RIGHT_ARROW,
                      x + Inches(3.85), Inches(1.65), Inches(0.3), Inches(0.3), ESDC_TEAL)

    # Remaining actions callout
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.4),
              Inches(12.3), Inches(0.04), ESDC_CORAL)

    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
                "IMMEDIATE ACTIONS:  ① Ouray: Review R&Rs on SharePoint, provide final comments  "
                "② Inder: Schedule System Team enablement session  "
                "③ Inder: Schedule stakeholder alignment with Chris and Don",
                font_size=12, font_color=ESDC_DARK, bold=False)

    add_slide_number(slide, n)
    add_speaker_notes(slide,
        "TALKING POINTS:\n"
        "We're using a proven 3-step approach: Align, Activate, Embed.\n\n"
        "ALIGN: We finalize the Charter — purpose, scope, R&Rs, decision rights. We also document "
        "what changes vs. what stays the same so the team knows exactly what's different.\n\n"
        "ACTIVATE: We run a focused enablement session covering the QPA learnings and Charter walkthrough. "
        "This is where we get ratification — not just awareness.\n\n"
        "EMBED: Over the next two sprints, we operationalize the Charter in daily ceremonies. "
        "We want to see it referenced in standups and refinement, not sitting on SharePoint.\n\n"
        "Three immediate actions are at the bottom — Ouray on R&R review, Inder on scheduling.\n\n"
        "ANTICIPATED Q: 'What does QPA stand for and why does it matter here?'\n"
        "A: Quality Process Assessment — it gave us data-driven insights on where our process gaps are, "
        "and those findings directly inform what the Charter needs to address.\n\n"
        "TRANSITION: Let me show you the milestones.")


def build_charter_milestones(prs):
    """Slide 7: Charter — Closure milestones."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    n = 7

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "System Team Charter", font_size=14, font_color=ESDC_GREY)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.5),
                "Four Auditable Milestones Track Charter from Draft to Operating Rhythm",
                font_size=22, font_color=ESDC_DARK, bold=True)

    milestones = [
        {
            "id": "M1", "title": "Charter v1.0 Ready for Ratification",
            "owner": "Inder (with Ouray review)",
            "exit": "Comments resolved, final version posted, ratification checklist complete"
        },
        {
            "id": "M2", "title": "System Team Activated",
            "owner": "Inder",
            "exit": "Enablement session complete, attendance captured, QPA & Charter deck published"
        },
        {
            "id": "M3", "title": "Stakeholder Alignment Completed",
            "owner": "Inder",
            "exit": "Chris & Don aligned on CI priorities; explicit start/stop/continue decisions captured"
        },
        {
            "id": "M4", "title": "Embedded into Operating Rhythm",
            "owner": "System Team Lead/SM (Inder support)",
            "exit": "Charter referenced in cadence, escalations follow defined path, first inspection done"
        },
    ]

    for i, ms in enumerate(milestones):
        y = Inches(1.6) + Inches(i * 1.3)

        # Milestone badge
        badge = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.5), y, Inches(0.8), Inches(0.5), ESDC_TEAL)
        p = badge.text_frame.paragraphs[0]
        p.text = ms["id"]
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Title
        add_textbox(slide, Inches(1.5), y, Inches(6), Inches(0.4),
                    ms["title"], font_size=16, font_color=ESDC_DARK, bold=True)

        # Owner
        add_textbox(slide, Inches(1.5), y + Inches(0.4), Inches(5), Inches(0.3),
                    f"Owner: {ms['owner']}", font_size=12, font_color=ESDC_TEAL, bold=True)

        # Exit criteria
        add_textbox(slide, Inches(1.5), y + Inches(0.7), Inches(10), Inches(0.4),
                    f"Exit: {ms['exit']}", font_size=12, font_color=ESDC_GREY)

        # Connecting line
        if i < 3:
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      Inches(0.85), y + Inches(1.05), Inches(0.04), Inches(0.25), ESDC_TEAL_LIGHT)

    # Value box
    val_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7), ESDC_LIGHT_GREY)
    val_box.line.color.rgb = ESDC_TEAL
    val_box.line.width = Pt(1)

    add_textbox(slide, Inches(0.7), Inches(6.35), Inches(2), Inches(0.3),
                "LEADERSHIP VALUE:", font_size=12, font_color=ESDC_TEAL, bold=True)
    add_textbox(slide, Inches(0.7), Inches(6.6), Inches(11.8), Inches(0.35),
                "A ratified operating model reduces ambiguity, speeds decision-making, "
                "and improves cross-team execution by making ownership and escalation paths explicit.",
                font_size=13, font_color=ESDC_DARK)

    add_slide_number(slide, n)
    add_speaker_notes(slide,
        "TALKING POINTS:\n"
        "Four auditable milestones — each with a clear owner and exit criteria.\n\n"
        "M1 is imminent — Ouray's review is the critical path. M2 follows with the enablement session. "
        "M3 is the Chris/Don alignment, which is one of my asks today. M4 is the proof point — "
        "are we actually using this in our ceremonies?\n\n"
        "The leadership value at the bottom is the 'so what': less ambiguity, faster decisions, "
        "better cross-team execution.\n\n"
        "TRANSITION: Now let's look at how we measure all of this — the Quality Dashboard.")


def build_dashboard_section_header(prs):
    """Slide 8: Section header — Quality Dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_divider_accent(slide)
    n = 8

    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(8), Inches(0.6),
                "Quality Dashboard & KPIs", font_size=32, font_color=ESDC_DARK, bold=True)

    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(0.5),
                "Deliver a dashboard leaders can trust for release health, shift-left progress, and intervention",
                font_size=16, font_color=ESDC_GREY)

    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(3.0), Inches(0.05), ESDC_CORAL)

    add_slide_number(slide, n)
    add_speaker_notes(slide, "TRANSITION: Section 2 — Quality Dashboard & KPIs.")


def build_dashboard_detail(prs):
    """Slide 9: Dashboard — closure plan and risk."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    n = 9

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "Quality Dashboard & KPIs", font_size=14, font_color=ESDC_GREY)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.5),
                "Dashboard Delivery Is Green with One Managed Dependency on Power BI Activation",
                font_size=22, font_color=ESDC_DARK, bold=True)

    add_status_badge(slide, Inches(11.3), Inches(0.78), "GREEN", GREEN)

    # Three-column closure plan
    phases = [
        {
            "label": "A",
            "title": "Finish CI Board Build",
            "color": ESDC_TEAL,
            "items": [
                "Lock baseline KPI definitions:\none definition, one calculation, one source",
                "Integrate RCA KPI with Ramesh:\ndefinition, logic, ownership, refresh cadence",
            ],
            "exit": "CI board shows baseline + RCA KPI, ready for leadership consumption"
        },
        {
            "label": "B",
            "title": "Unlock Power BI Activation",
            "color": ESDC_CORAL,
            "items": [
                "RISK: Without Power BI, cannot reliably\ntrend time-series across sprints/releases",
                "MITIGATION: Manual calc + screenshots\npreserve trend continuity in weekly deck",
                "ACCELERATOR: Short-term QE backfill\nto unblock KPI automation (fastest path)",
            ],
            "exit": "Power BI layer active; time-series trending automated"
        },
        {
            "label": "C",
            "title": "Embed into Governance",
            "color": ESDC_TEAL_DARK,
            "items": [
                "Default view for weekly release health",
                "CI prioritization discussion input",
                "Executive updates: trend, not snapshot",
            ],
            "exit": "Dashboard = single source of truth for CI signals and weekly reporting"
        },
    ]

    for i, ph in enumerate(phases):
        x = Inches(0.5) + Inches(i * 4.2)

        # Phase header
        hdr_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                            x, Inches(1.55), Inches(3.8), Inches(0.5), ph["color"])
        p = hdr_box.text_frame.paragraphs[0]
        p.text = f"{ph['label']}. {ph['title']}"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Bullet items
        add_bullet_frame(slide, x + Inches(0.1), Inches(2.3), Inches(3.6), Inches(2.8),
                         ph["items"], font_size=12, bullet_color=ph["color"], spacing_after=Pt(8))

        # Exit box
        ex_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                           x, Inches(5.1), Inches(3.8), Inches(0.8), ESDC_LIGHT_GREY)
        ex_box.line.color.rgb = ph["color"]
        ex_box.line.width = Pt(1)

        add_textbox(slide, x + Inches(0.1), Inches(5.15), Inches(3.6), Inches(0.2),
                    "EXIT CRITERIA:", font_size=10, font_color=ph["color"], bold=True)
        add_textbox(slide, x + Inches(0.1), Inches(5.35), Inches(3.6), Inches(0.5),
                    ph["exit"], font_size=11, font_color=ESDC_DARK)

    # Value box
    val_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7), ESDC_LIGHT_GREY)
    val_box.line.color.rgb = ESDC_TEAL
    val_box.line.width = Pt(1)

    add_textbox(slide, Inches(0.7), Inches(6.35), Inches(2), Inches(0.3),
                "LEADERSHIP VALUE:", font_size=12, font_color=ESDC_TEAL, bold=True)
    add_textbox(slide, Inches(0.7), Inches(6.6), Inches(11.8), Inches(0.35),
                "Moves decisions from opinion-based to signal-based. Enables early intervention, "
                "tracks shift-left progress, and makes continuous improvement measurable over time.",
                font_size=13, font_color=ESDC_DARK)

    add_slide_number(slide, n)
    add_speaker_notes(slide,
        "TALKING POINTS:\n"
        "Overall status is Green with one managed dependency — Power BI activation.\n\n"
        "Three-part closure plan:\n"
        "A) We're finalizing the CI board build. One definition, one calculation, one source — "
        "this eliminates ambiguity about what each KPI means.\n"
        "B) Power BI is a risk I want to flag. Without it, we can't automate time-series trending. "
        "We have a mitigation in place (manual screenshots), but the fastest path to resolution is "
        "a short-term QE backfill. This is one of my asks.\n"
        "C) Once built, the dashboard becomes THE view for weekly release health — not a separate artifact.\n\n"
        "ANTICIPATED Q: 'What's the dependency on Tim Miller's team?'\n"
        "A: Power BI activation requires capacity from Tim's team. We've logged the dependency "
        "and have the manual mitigation, but automation is the sustainable path.\n\n"
        "ANTICIPATED Q: 'How long would the QE backfill be needed?'\n"
        "A: Estimated 2-3 sprints of dedicated support to build the Power BI layer.\n\n"
        "TRANSITION: Now let's look at how we activate enablement through the ARTs and STEs.")


def build_enablement_section_header(prs):
    """Slide 10: Section header — Enablement Activation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_section_divider_accent(slide)
    n = 10

    add_textbox(slide, Inches(0.5), Inches(1.0), Inches(9), Inches(0.6),
                "Enablement Activation", font_size=32, font_color=ESDC_DARK, bold=True)

    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(0.5),
                "Accessibility & Automation operationalized through ARTs and STEs, "
                "tracked via Quality Dashboard",
                font_size=16, font_color=ESDC_GREY)

    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.5), Inches(3.0), Inches(0.05), ESDC_TEAL)

    add_slide_number(slide, n)
    add_speaker_notes(slide, "TRANSITION: Section 3 — Enablement Activation.")


def build_enablement_detail(prs):
    """Slide 11: Enablement — activation sequence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    n = 11

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "Enablement Activation", font_size=14, font_color=ESDC_GREY)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.5),
                "Three Meetings Move Enablement from Plan to Product Team Commitments",
                font_size=22, font_color=ESDC_DARK, bold=True)

    meetings = [
        {
            "num": "1", "title": "ALIGN\nInternal Greenlight",
            "date": "Thu Feb 12",
            "duration": "45 min",
            "color": ESDC_TEAL_LIGHT,
            "attendees": "QE core team: Ouray, Venu, Yash,\ndashboard owner, enablement leads",
            "decisions": [
                "Approve activation plan v1.0",
                "Confirm RACI and schedule",
                "Confirm activation kit & escalation path",
            ]
        },
        {
            "num": "2", "title": "ACTIVATE\nARTs & STEs Kickoff",
            "date": "Wed Feb 18",
            "duration": "60 min",
            "color": ESDC_TEAL,
            "attendees": "Intake ART (Jacqueline), Manage ART\n(Laura), Platform ART, Mike Cattran,\nChris Snow, STEs, QE support",
            "decisions": [
                "Ratify cadence & 30/60/90 measures",
                "Confirm roles",
                "Approve product team activation approach",
            ]
        },
        {
            "num": "3", "title": "EMBED\nProduct Team Activation",
            "date": "Week of Feb 23",
            "duration": "30 min/team",
            "color": ESDC_TEAL_DARK,
            "attendees": "PO, SM, Dev Lead, QE Lead,\nA11y & Automation champions,\nSTEs, ART rep (optional), Inder",
            "decisions": [
                "Lock 2–3 commitments with owners & dates",
                "Confirm tracking cadence",
                "Log data quality/tooling fixes",
            ]
        },
    ]

    for i, mtg in enumerate(meetings):
        x = Inches(0.5) + Inches(i * 4.2)

        # Header with meeting number and title
        hdr = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                        x, Inches(1.55), Inches(3.8), Inches(0.85), mtg["color"])
        tf = hdr.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = mtg["title"]
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE if mtg["color"] != ESDC_TEAL_LIGHT else ESDC_TEAL_DARK
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Date & duration
        add_textbox(slide, x, Inches(2.55), Inches(3.8), Inches(0.35),
                    f"📅 {mtg['date']}  •  ⏱ {mtg['duration']}",
                    font_size=12, font_color=ESDC_TEAL_DARK, bold=True, alignment=PP_ALIGN.CENTER)

        # Attendees
        add_textbox(slide, x + Inches(0.1), Inches(2.95), Inches(3.6), Inches(0.2),
                    "ATTENDEES:", font_size=10, font_color=ESDC_GREY, bold=True)
        add_textbox(slide, x + Inches(0.1), Inches(3.15), Inches(3.6), Inches(0.9),
                    mtg["attendees"], font_size=11, font_color=ESDC_DARK)

        # Decisions
        add_textbox(slide, x + Inches(0.1), Inches(4.1), Inches(3.6), Inches(0.2),
                    "DECISIONS & OUTPUTS:", font_size=10, font_color=mtg["color"]
                    if mtg["color"] != ESDC_TEAL_LIGHT else ESDC_TEAL, bold=True)

        add_bullet_frame(slide, x + Inches(0.1), Inches(4.35), Inches(3.6), Inches(1.5),
                         mtg["decisions"], font_size=12,
                         bullet_color=mtg["color"] if mtg["color"] != ESDC_TEAL_LIGHT else ESDC_TEAL,
                         spacing_after=Pt(6))

        # Arrow connector
        if i < 2:
            add_shape(slide, MSO_SHAPE.RIGHT_ARROW,
                      x + Inches(3.85), Inches(1.85), Inches(0.3), Inches(0.3), ESDC_TEAL)

    # Value box
    val_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7), ESDC_LIGHT_GREY)
    val_box.line.color.rgb = ESDC_TEAL
    val_box.line.width = Pt(1)

    add_textbox(slide, Inches(0.7), Inches(6.35), Inches(2), Inches(0.3),
                "LEADERSHIP VALUE:", font_size=12, font_color=ESDC_TEAL, bold=True)
    add_textbox(slide, Inches(0.7), Inches(6.6), Inches(11.8), Inches(0.35),
                "Operationalizes the System Team Charter by creating a consistent, scalable interface "
                "where the System Team enables Product Teams through ART/STE channels, using metrics as "
                "the accountability mechanism.",
                font_size=13, font_color=ESDC_DARK)

    add_slide_number(slide, n)
    add_speaker_notes(slide,
        "TALKING POINTS:\n"
        "Three meetings, three outcomes — each building on the last.\n\n"
        "Meeting 1 (Feb 12): Internal greenlight. We align the core QE team on scope, "
        "sequencing, and the activation kit before going external.\n\n"
        "Meeting 2 (Feb 18): ARTs and STEs kickoff. This is where we ratify the engagement model "
        "and 30/60/90 success measures with the actual ART leads — Jacqueline, Laura, and others.\n\n"
        "Meeting 3 (Week of Feb 23): Product team sessions. This is where the rubber meets the road. "
        "Each team locks 2-3 specific commitments with owners and dates. We use their team dashboards "
        "as the conversation starter.\n\n"
        "ANTICIPATED Q: 'Which product teams go first?'\n"
        "A: Focus teams first — those with the highest need or readiness. We'll confirm sequencing "
        "in the internal greenlight session.\n\n"
        "TRANSITION: Let me pull this together into a consolidated timeline.")


def build_timeline(prs):
    """Slide 12: Consolidated timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    n = 12

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "EI on BDM  —  Leadership Update", font_size=14, font_color=ESDC_GREY)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.5),
                "Consolidated Timeline: All Three Workstreams Converge by End of February",
                font_size=22, font_color=ESDC_DARK, bold=True)

    # Timeline header row
    weeks = ["Feb 10–14", "Feb 17–21", "Feb 24–28", "Mar 3+"]
    labels = ["THIS WEEK", "WEEK 2", "WEEK 3", "ONGOING"]

    for i, (wk, lbl) in enumerate(zip(weeks, labels)):
        x = Inches(3.5) + Inches(i * 2.5)
        add_callout_box(slide, x, Inches(1.5), Inches(2.2), Inches(0.55),
                        f"{lbl}\n{wk}",
                        fill=ESDC_TEAL if i < 3 else ESDC_GREY, font_size=11)

    # Workstream rows
    streams = [
        {
            "name": "System Team\nCharter",
            "color": ESDC_TEAL,
            "items": [
                "M1: Charter\nready for\nratification",
                "M2: Enablement\nsession\n(Activate)",
                "M3: Chris/Don\nalignment",
                "M4: Embedded\nin rhythm\n(ongoing)"
            ]
        },
        {
            "name": "Quality\nDashboard",
            "color": ESDC_CORAL,
            "items": [
                "CI board +\nRCA KPI\nfinalized",
                "Power BI\nactivation\n(if unblocked)",
                "Dashboard\nembedded in\ngovernance",
                "Trend-based\nexecutive\nreporting"
            ]
        },
        {
            "name": "Enablement\nActivation",
            "color": ESDC_TEAL_DARK,
            "items": [
                "Internal\ngreenlight\n(Feb 12)",
                "ART/STE\nkickoff\n(Feb 18)",
                "Product team\nsessions\n(Feb 23+)",
                "Track\ncommitments\n& adoption"
            ]
        },
    ]

    for si, stream in enumerate(streams):
        y = Inches(2.35) + Inches(si * 1.55)

        # Stream label
        lbl_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.5), y, Inches(2.7), Inches(1.2), stream["color"])
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stream["name"]
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Weekly items
        for wi, item in enumerate(stream["items"]):
            x = Inches(3.5) + Inches(wi * 2.5)
            item_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                 x, y, Inches(2.2), Inches(1.2), ESDC_LIGHT_GREY)
            item_box.line.color.rgb = stream["color"]
            item_box.line.width = Pt(1)
            tf2 = item_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = item
            p2.font.size = Pt(11)
            p2.font.color.rgb = ESDC_DARK
            p2.alignment = PP_ALIGN.CENTER

    # "Now" marker
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(2.15),
              Inches(0.04), Inches(5.0), ESDC_CORAL)
    add_textbox(slide, Inches(3.1), Inches(2.0), Inches(0.8), Inches(0.25),
                "▼ NOW", font_size=10, font_color=ESDC_CORAL, bold=True)

    add_slide_number(slide, n)
    add_speaker_notes(slide,
        "TALKING POINTS:\n"
        "Here's the consolidated view. All three workstreams converge by end of February.\n\n"
        "This week is heavy: Charter M1 ready, CI board finalized, and the internal greenlight "
        "for enablement activation.\n\n"
        "Week 2: Enablement session for the Charter, ART/STE kickoff, and Power BI activation "
        "(if we get the capacity decision).\n\n"
        "Week 3: Product team activation sessions begin, dashboard embedded in governance.\n\n"
        "March onward is sustained operations — these become part of how we work, not projects.\n\n"
        "TRANSITION: Let me close with the three specific asks I have for you.")


def build_leadership_asks(prs):
    """Slide 13: Leadership asks and decisions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    n = 13

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "EI on BDM  —  Leadership Update", font_size=14, font_color=ESDC_GREY)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.5),
                "Three Decisions Needed to Maintain Momentum",
                font_size=26, font_color=ESDC_DARK, bold=True)

    asks = [
        {
            "num": "1",
            "title": "Confirm Sponsorship for Charter Activation & CI Alignment",
            "detail": "Schedule stakeholder alignment with HT Chris and Don on the Continuous Improvement "
                      "plan and prioritization to secure clear decisions and sponsorship.",
            "impact": "Without sponsorship, the Charter lacks organizational authority and CI priorities remain unresolved.",
            "color": ESDC_TEAL,
        },
        {
            "num": "2",
            "title": "Confirm Power BI Capacity or Approve QE Backfill",
            "detail": "Either confirm capacity from Tim Miller's Power BI team, or approve short-term dedicated "
                      "QE backfill to build the Power BI layer and unblock KPI automation.",
            "impact": "Without this, we cannot automate time-series trending and remain dependent on manual screenshots.",
            "color": ESDC_CORAL,
        },
        {
            "num": "3",
            "title": "Confirm Quality Dashboard as Default Executive View",
            "detail": "Confirm that the Quality Dashboard becomes the default view for weekly release "
                      "health checkpoints and executive updates going forward.",
            "impact": "Establishes a single source of truth and eliminates parallel reporting artifacts.",
            "color": ESDC_TEAL_DARK,
        },
    ]

    for i, ask in enumerate(asks):
        y = Inches(1.55) + Inches(i * 1.7)

        # Number badge
        badge = add_shape(slide, MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.65), Inches(0.65), ask["color"])
        p = badge.text_frame.paragraphs[0]
        p.text = ask["num"]
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Title
        add_textbox(slide, Inches(1.4), y, Inches(8), Inches(0.4),
                    ask["title"], font_size=17, font_color=ESDC_DARK, bold=True)

        # Detail
        add_textbox(slide, Inches(1.4), y + Inches(0.45), Inches(8), Inches(0.5),
                    ask["detail"], font_size=13, font_color=ESDC_DARK)

        # Impact callout
        impact_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(1.4), y + Inches(1.0), Inches(8), Inches(0.4), ESDC_LIGHT_GREY)
        impact_box.line.color.rgb = ask["color"]
        impact_box.line.width = Pt(1)

        tf = impact_box.text_frame
        tf.word_wrap = True
        p2 = tf.paragraphs[0]
        p2.text = f"⚡ WHY IT MATTERS:  {ask['impact']}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = ESDC_GREY
        p2.font.name = "Calibri"

    add_slide_number(slide, n)
    add_speaker_notes(slide,
        "TALKING POINTS:\n"
        "Three specific, time-bound decisions:\n\n"
        "Ask 1 — Sponsorship: I need the greenlight to schedule alignment with Chris and Don. "
        "The Charter needs organizational authority, not just team agreement.\n\n"
        "Ask 2 — Power BI: This is the most tangible risk. Two options: Tim's team provides capacity, "
        "or we approve a short-term QE backfill. I recommend the backfill as the faster path.\n\n"
        "Ask 3 — Dashboard as default: This is a mindset shift. One view for release health. "
        "No more separate decks or spreadsheets.\n\n"
        "ANTICIPATED Q: 'What's the cost of the QE backfill?'\n"
        "A: Estimated 2-3 sprints of one dedicated resource. The ROI is automated KPI trending "
        "that eliminates weekly manual effort permanently.\n\n"
        "ANTICIPATED Q: 'Are Chris and Don aware this is coming?'\n"
        "A: They have context from previous discussions, but we need a formal alignment to secure "
        "explicit decisions on CI priorities.\n\n"
        "CLOSE: I'm looking for your direction on these three items today if possible.")


def build_closing(prs):
    """Slide 14: Closing / call to action."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    n = 14

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "EI on BDM  —  Leadership Update", font_size=14, font_color=ESDC_GREY)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.5),
                "Summary: From Process to Practice — Closing the Loop",
                font_size=24, font_color=ESDC_DARK, bold=True)

    # Three summary boxes
    summaries = [
        ("System Team Charter", "Operating model\nratified & activated", ESDC_TEAL, "Reduces ambiguity,\nspeeds decisions"),
        ("Quality Dashboard", "Signal-based\ndecision making", ESDC_CORAL, "Early intervention,\nmeasurable CI"),
        ("Enablement Activation", "Product teams\ncommitted & tracked", ESDC_TEAL_DARK, "Scalable, consistent\ninterface"),
    ]

    for i, (title, what, color, value) in enumerate(summaries):
        x = Inches(0.5) + Inches(i * 4.2)

        add_callout_box(slide, x, Inches(1.6), Inches(3.8), Inches(0.5),
                        title, fill=color, font_size=14)

        # What
        add_textbox(slide, x + Inches(0.3), Inches(2.3), Inches(3.2), Inches(0.8),
                    what, font_size=16, font_color=ESDC_DARK, bold=True, alignment=PP_ALIGN.CENTER)

        # Value
        add_textbox(slide, x + Inches(0.3), Inches(3.2), Inches(3.2), Inches(0.8),
                    value, font_size=13, font_color=ESDC_GREY, alignment=PP_ALIGN.CENTER)

    # Convergence statement
    conv_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                         Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.8), ESDC_TEAL)

    tf = conv_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("Together, these three enablers create a closed-loop system: "
              "the Charter defines HOW we operate, the Dashboard shows HOW WELL we're doing, "
              "and Enablement puts the tools in teams' hands.")
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = False
    p.alignment = PP_ALIGN.CENTER

    # Next steps
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(4), Inches(0.35),
                "IMMEDIATE NEXT STEPS", font_size=15, font_color=ESDC_DARK, bold=True)

    next_steps = [
        "Ouray: Complete R&R review (this week)",
        "Inder: Schedule enablement session and Chris/Don alignment",
        "Internal greenlight meeting: February 12",
        "ART/STE kickoff: February 18",
        "Product team activation: week of February 23",
    ]
    add_bullet_frame(slide, Inches(0.5), Inches(5.9), Inches(12), Inches(1.3),
                     next_steps, font_size=13, bullet_color=ESDC_TEAL, spacing_after=Pt(5))

    add_slide_number(slide, n)
    add_speaker_notes(slide,
        "CLOSING:\n"
        "To summarize: we're closing the loop. The Charter, the Dashboard, and Enablement Activation "
        "aren't three separate initiatives — they're one system.\n\n"
        "The next steps are clear, the owners are identified, and the timeline is tight but achievable.\n\n"
        "I need your direction on the three asks to keep this on track.\n\n"
        "Thank you. I'm happy to take questions.\n\n"
        "BACKUP READY:\n"
        "- Detailed RACI for enablement activation\n"
        "- KPI definitions and calculation logic\n"
        "- Full Charter v1.0 draft\n"
        "- QPA findings summary")


def build_appendix_raci(prs):
    """Slide 15: Appendix — Enablement Activation RACI (backup)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_esdc_header_bar(slide)
    n = 15

    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(10), Inches(0.45),
                "APPENDIX", font_size=14, font_color=ESDC_CORAL, bold=True)

    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.5),
                "Enablement Activation — Detailed Participant Map",
                font_size=22, font_color=ESDC_DARK, bold=True)

    # Table header
    headers = ["Meeting", "Purpose", "Key Attendees", "Key Decisions"]
    col_widths = [Inches(2.0), Inches(3.0), Inches(4.0), Inches(3.5)]

    y_start = Inches(1.5)
    row_h = Inches(0.4)
    x_pos = Inches(0.3)

    # Header row
    for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
        cell = add_shape(slide, MSO_SHAPE.RECTANGLE, x_pos, y_start, cw, row_h, ESDC_TEAL)
        p = cell.text_frame.paragraphs[0]
        p.text = hdr
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        x_pos += cw

    rows = [
        ["1. Internal\nGreenlight\n(Feb 12, 45m)",
         "Align scope, sequencing,\nroles, artifacts, comms,\nsuccess measures",
         "Ouray, Venu, Yash,\ndashboard/metrics owner,\nenablement leads",
         "Approve plan v1.0,\nconfirm RACI/schedule,\nconfirm escalation path"],
        ["2. ART/STE\nKickoff\n(Feb 18, 60m)",
         "Ratify enablement model\nand expectations, confirm\ncadence & engagement",
         "Jacqueline (Intake ART),\nLaura (Manage ART),\nPlatform ART, Mike Cattran,\nChris Snow, STEs, QE",
         "Ratify cadence,\n30/60/90 measures,\nconfirm roles, approve\nPT activation approach"],
        ["3. Product Team\nActivation\n(Feb 23+, 30m/team)",
         "Use dashboards & metrics\nto secure adoption\ncommitments",
         "PO, SM, Dev Lead,\nQE Lead, A11y/Automation\nchampions, STEs, ART rep,\nInder (facilitating)",
         "Lock 2–3 commitments\nw/ owners & dates,\nconfirm tracking cadence,\nlog tooling fixes"],
    ]

    for ri, row in enumerate(rows):
        y = y_start + row_h + Inches(ri * 1.55)
        x_pos = Inches(0.3)
        bg = WHITE if ri % 2 == 0 else ESDC_LIGHT_GREY
        for j, (cell_text, cw) in enumerate(zip(row, col_widths)):
            cell = add_shape(slide, MSO_SHAPE.RECTANGLE, x_pos, y, cw, Inches(1.4), bg)
            cell.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            cell.line.width = Pt(0.5)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell_text
            p.font.size = Pt(10)
            p.font.color.rgb = ESDC_DARK
            p.font.name = "Calibri"
            if j == 0:
                p.font.bold = True
            x_pos += cw

    add_slide_number(slide, n)
    add_speaker_notes(slide,
        "BACKUP SLIDE:\n"
        "This is the detailed participant map for the enablement activation sequence. "
        "Use this if leadership asks for specifics on who's involved and what each meeting decides.")


# ── Main Build ─────────────────────────────────────────────────────────

def main():
    if TEMPLATE_PATH and os.path.exists(TEMPLATE_PATH):
        prs = Presentation(TEMPLATE_PATH)
        print(f"✅ Loaded template: {TEMPLATE_PATH}")
    else:
        prs = Presentation()
        prs.slide_width = Emu(12192000)   # 13.333 inches (widescreen)
        prs.slide_height = Emu(6858000)   # 7.5 inches
        print("ℹ️  No template found — building from ESDC brand guidelines.")

    # Build all slides
    build_title_slide(prs)           # Slide 1
    build_exec_summary(prs)          # Slide 2
    build_agenda(prs)                # Slide 3
    build_strategic_context(prs)     # Slide 4
    build_charter_section_header(prs)  # Slide 5
    build_charter_align_activate_embed(prs)  # Slide 6
    build_charter_milestones(prs)    # Slide 7
    build_dashboard_section_header(prs)  # Slide 8
    build_dashboard_detail(prs)      # Slide 9
    build_enablement_section_header(prs)  # Slide 10
    build_enablement_detail(prs)     # Slide 11
    build_timeline(prs)              # Slide 12
    build_leadership_asks(prs)       # Slide 13
    build_closing(prs)               # Slide 14
    build_appendix_raci(prs)         # Slide 15

    prs.save(OUTPUT_FILE)
    print(f"\n✅ Presentation saved: {OUTPUT_FILE}")
    print(f"   {len(prs.slides)} slides generated with speaker notes.")

if __name__ == "__main__":
    main()